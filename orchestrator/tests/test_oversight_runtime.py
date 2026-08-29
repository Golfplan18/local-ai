"""Tests for the corpus and output runtime modules.

Covers C-Instance, C-Validate, and O-Render against synthetic templates,
instances, and OFF specs. Also verifies the oversight events fire on each
runtime invocation.
"""
from __future__ import annotations

import os
import shutil
import sys
import tempfile
import unittest
from textwrap import dedent

HERE = os.path.dirname(os.path.abspath(__file__))
ORCH = os.path.dirname(HERE)
if ORCH not in sys.path:
    sys.path.insert(0, ORCH)
if HERE not in sys.path:
    sys.path.insert(0, HERE)

from corpus_runtime import c_instance, c_validate  # noqa: E402
from corpus_parser import parse_corpus_file  # noqa: E402
from output_runtime import o_render, parse_off_spec  # noqa: E402
from oversight_events import register_handler, clear_handlers  # noqa: E402
from oversight_sandbox import redirect_oversight_logs  # noqa: E402


SAMPLE_TEMPLATE = dedent("""\
    ---
    type: corpus_template
    template_version: 1.0
    ---

    # Marketing Monthly Corpus Template

    ## Sections

    ```yaml
    sections:
      - id: weekly_sales
        name: Weekly Sales
        source: pff-mortgage-pipeline
        missing_data_behavior: hold-and-warn
      - id: campaigns
        name: Campaign Performance
        source: pff-campaign-extractor
        missing_data_behavior: default-empty
    ```
    """)

SAMPLE_OFF = dedent("""\
    ---
    name: monthly-board-memo
    medium: markdown
    title: "Monthly Memo — {period}"
    sections:
      - section: weekly_sales
        heading: Weekly Sales
      - section: campaigns
        heading: Campaign Performance
    ---
    """)


class _TempWorkspace(unittest.TestCase):
    """Base class providing a clean tempdir per test."""

    def setUp(self):
        redirect_oversight_logs(self)
        self.tmp = tempfile.mkdtemp(prefix="ora-runtime-test-")
        self.template_path = os.path.join(self.tmp, "template.md")
        self.off_path = os.path.join(self.tmp, "off.md")
        self.instance_dir = os.path.join(self.tmp, "instances")
        self.output_dir = os.path.join(self.tmp, "outputs")
        with open(self.template_path, "w") as f:
            f.write(SAMPLE_TEMPLATE)
        with open(self.off_path, "w") as f:
            f.write(SAMPLE_OFF)
        clear_handlers()
        self.events = []
        register_handler(lambda e: self.events.append(e["event_type"]))

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)
        clear_handlers()


class TestCInstance(_TempWorkspace):

    def test_creates_instance_file(self):
        result = c_instance(self.template_path, "2026-05", self.instance_dir)
        self.assertTrue(result.success, result.error)
        self.assertTrue(os.path.isfile(result.instance_path))
        self.assertEqual(result.period, "2026-05")
        self.assertEqual(set(result.sections_created), {"weekly_sales", "campaigns"})

    def test_emits_corpus_instance_created_event(self):
        c_instance(self.template_path, "2026-05", self.instance_dir,
                   workflow_id="smoke", project_nexus="smoke_test")
        self.assertIn("CorpusInstanceCreated", self.events)

    def test_refuses_to_overwrite_existing_instance(self):
        c_instance(self.template_path, "2026-05", self.instance_dir)
        result2 = c_instance(self.template_path, "2026-05", self.instance_dir)
        self.assertFalse(result2.success)
        self.assertIn("already exists", result2.error)

    def test_handles_missing_template(self):
        result = c_instance("/nonexistent/template.md", "2026-05", self.instance_dir)
        self.assertFalse(result.success)
        self.assertIn("not found", result.error)

    def test_includes_section_headings_for_each_template_section(self):
        result = c_instance(self.template_path, "2026-05", self.instance_dir)
        with open(result.instance_path) as f:
            content = f.read()
        self.assertIn("### Section weekly_sales — Weekly Sales", content)
        self.assertIn("### Section campaigns — Campaign Performance", content)

    def test_rejects_malformed_template_frontmatter(self):
        with open(self.template_path, "w") as f:
            f.write("---\ntype: [corpus_template\n---\n# Broken\n")
        result = c_instance(self.template_path, "2026-05", self.instance_dir)
        self.assertFalse(result.success)
        self.assertIn("Invalid YAML frontmatter", result.error)
        self.assertFalse(os.path.exists(self.instance_dir))

    def test_rejects_falsey_non_mapping_frontmatter_and_invalid_tags(self):
        body = "# Corpus\n\n## Sections\n\n### Section SalesQ2 — Quarterly Sales\n"
        cases = {
            "empty-list-root": "[]",
            "false-root": "false",
            "zero-root": "0",
            "null-root": "null",
            "empty-mapping-tags": "type: incubator\ntags: {}",
            "false-tags": "type: incubator\ntags: false",
            "zero-tags": "type: incubator\ntags: 0",
            "null-tags": "type: incubator\ntags: null",
        }
        for label, frontmatter in cases.items():
            with self.subTest(case=label):
                with open(self.template_path, "w") as f:
                    f.write(f"---\n{frontmatter}\n---\n\n{body}")
                destination = os.path.join(self.tmp, f"instances-{label}")
                result = c_instance(self.template_path, "2026-05", destination)
                self.assertFalse(result.success)
                self.assertIn("Invalid YAML frontmatter", result.error)
                self.assertFalse(os.path.exists(destination))
                self.assertNotIn("CorpusInstanceCreated", self.events)

    def test_period_with_colon_round_trips_through_valid_yaml(self):
        period = "Q2 2026: revised"
        result = c_instance(self.template_path, period, self.instance_dir)
        self.assertTrue(result.success, result.error)
        parsed = parse_corpus_file(result.instance_path)
        self.assertTrue(parsed.is_valid, parsed.parse_error)
        self.assertFalse(parsed.is_template)
        self.assertEqual(parsed.instance_period, period)


class TestCValidate(_TempWorkspace):

    def setUp(self):
        super().setUp()
        # Create a fresh instance to validate
        result = c_instance(self.template_path, "2026-05", self.instance_dir)
        self.assertTrue(result.success)
        self.instance_path = result.instance_path

    def test_unpopulated_instance_returns_fail(self):
        result = c_validate(self.instance_path, self.template_path)
        self.assertEqual(result.overall_status, "FAIL")

    def test_partially_populated_returns_partial(self):
        # Populate only weekly_sales
        with open(self.instance_path) as f:
            content = f.read()
        content = content.replace(
            "<!-- Section content goes here. -->",
            "Sales data here.",
            1,  # only first occurrence
        )
        with open(self.instance_path, "w") as f:
            f.write(content)

        result = c_validate(self.instance_path, self.template_path)
        # weekly_sales populated, campaigns empty (default-empty acceptable -> still validates)
        # Actual behavior: 1/2 populated, no missing → PARTIAL
        self.assertIn(result.overall_status, ("PARTIAL", "PASS"))
        self.assertEqual(result.populated_count, 1)

    def test_fully_populated_returns_pass(self):
        with open(self.instance_path) as f:
            content = f.read()
        content = content.replace(
            "<!-- Section content goes here. -->",
            "Some real content here.",
        )
        with open(self.instance_path, "w") as f:
            f.write(content)
        result = c_validate(self.instance_path, self.template_path)
        self.assertEqual(result.overall_status, "PASS")
        self.assertTrue(result.success)

    def test_emits_corpus_validated_event(self):
        c_validate(self.instance_path, self.template_path,
                   workflow_id="smoke", project_nexus="smoke_test")
        self.assertIn("CorpusValidated", self.events)

    def test_handles_missing_instance(self):
        result = c_validate("/nonexistent/instance.md")
        self.assertEqual(result.overall_status, "FAIL")
        self.assertIn("not found", result.error)

    def test_rejects_malformed_instance_frontmatter(self):
        with open(self.instance_path, "w") as f:
            f.write("---\ntype: [corpus_instance\n---\n# Broken\n")
        result = c_validate(self.instance_path, self.template_path)
        self.assertFalse(result.success)
        self.assertEqual(result.overall_status, "FAIL")
        self.assertIn("Invalid YAML frontmatter", result.error)


class TestORender(_TempWorkspace):

    def setUp(self):
        super().setUp()
        result = c_instance(self.template_path, "2026-05", self.instance_dir)
        self.instance_path = result.instance_path
        # Populate both sections
        with open(self.instance_path) as f:
            content = f.read()
        content = content.replace(
            "<!-- Section content goes here. -->",
            "Sample content here.",
        )
        with open(self.instance_path, "w") as f:
            f.write(content)

    def test_renders_markdown_artifact(self):
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        self.assertTrue(result.success, result.error)
        self.assertTrue(os.path.isfile(result.artifact_path))
        self.assertEqual(result.artifact_format, "markdown")

    def test_artifact_contains_section_headings(self):
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        with open(result.artifact_path) as f:
            content = f.read()
        self.assertIn("## Weekly Sales", content)
        self.assertIn("## Campaign Performance", content)

    def test_emits_off_rendered_event(self):
        o_render(self.off_path, self.instance_path, self.output_dir,
                 workflow_id="smoke", project_nexus="smoke_test")
        self.assertIn("OFFRendered", self.events)

    def test_title_interpolation(self):
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        with open(result.artifact_path) as f:
            content = f.read()
        self.assertIn("Monthly Memo — 2026-05", content)

    def test_handles_missing_off_spec(self):
        result = o_render("/nonexistent/off.md", self.instance_path, self.output_dir)
        self.assertFalse(result.success)
        self.assertIn("not found", result.error)

    def test_rejects_malformed_corpus_frontmatter_without_artifact_or_event(self):
        with open(self.instance_path, "w") as f:
            f.write("---\ntype: [corpus_instance\n---\n# Broken\n")
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        self.assertFalse(result.success)
        self.assertIn("Invalid YAML frontmatter", result.error)
        self.assertFalse(os.path.exists(self.output_dir))
        self.assertNotIn("OFFRendered", self.events)

    def test_rejects_malformed_or_structureless_off_before_output_or_event(self):
        cases = {
            "malformed-yaml": "---\nname: broken\nsections: [\n---\n",
            "missing-sections": dedent("""\
                ---
                name: no-sections
                medium: markdown
                ---

                # This is not a renderable OFF
                """),
        }
        for label, off_text in cases.items():
            with self.subTest(case=label):
                off_path = os.path.join(self.tmp, f"{label}.md")
                output_dir = os.path.join(self.tmp, f"outputs-{label}")
                with open(off_path, "w") as f:
                    f.write(off_text)
                self.events.clear()

                result = o_render(off_path, self.instance_path, output_dir)

                self.assertFalse(result.success)
                self.assertIn("Failed to parse OFF spec", result.error)
                self.assertFalse(os.path.exists(output_dir))
                self.assertNotIn("OFFRendered", self.events)

    def test_valid_minimal_body_marker_off_still_renders(self):
        minimal_off = os.path.join(self.tmp, "minimal-body-off.md")
        with open(minimal_off, "w") as f:
            f.write("## Section: weekly_sales\n")

        result = o_render(minimal_off, self.instance_path, self.output_dir)

        self.assertTrue(result.success, result.error)
        self.assertTrue(os.path.isfile(result.artifact_path))
        self.assertEqual(result.sections_read, ["weekly_sales"])
        self.assertIn("OFFRendered", self.events)

    def test_escaped_pipe_table_has_identical_artifact_cells_in_all_formats(self):
        with open(self.instance_path) as f:
            content = f.read()
        content = content.replace(
            "Sample content here.",
            (
                "| Label | Detail |\n"
                "| --- | --- |\n"
                "| Qualified \\| converted | Retained |"
            ),
            1,
        )
        with open(self.instance_path, "w") as f:
            f.write(content)

        artifacts = {}
        for medium in ("html", "docx", "pptx", "xlsx"):
            off_path = os.path.join(self.tmp, f"off-{medium}.md")
            with open(off_path, "w") as f:
                f.write(dedent(f"""\
                    ---
                    name: escaped-pipe-{medium}
                    medium: {medium}
                    sections:
                      - section: weekly_sales
                        heading: Metrics
                    ---
                    """))
            result = o_render(off_path, self.instance_path, self.output_dir)
            self.assertTrue(result.success, result.error)
            artifacts[medium] = result.artifact_path

        import html as _html
        import re as _re
        with open(artifacts["html"]) as f:
            html_text = f.read()
        html_rows = []
        for row_html in _re.findall(r"<tr>(.*?)</tr>", html_text, _re.DOTALL):
            cells = _re.findall(r"<t[hd]>(.*?)</t[hd]>", row_html, _re.DOTALL)
            html_rows.append([_html.unescape(_re.sub(r"<[^>]+>", "", c)) for c in cells])

        from docx import Document
        doc = Document(artifacts["docx"])
        docx_rows = [[cell.text for cell in row.cells] for row in doc.tables[0].rows]

        from pptx import Presentation
        prs = Presentation(artifacts["pptx"])
        pptx_table = next(
            shape.table
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_table
        )
        pptx_rows = [[cell.text for cell in row.cells] for row in pptx_table.rows]

        import xml.etree.ElementTree as ET
        import zipfile
        spreadsheet_ns = {"x": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        with zipfile.ZipFile(artifacts["xlsx"]) as archive:
            shared_root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
            shared_strings = [
                "".join(node.text or "" for node in item.findall(".//x:t", spreadsheet_ns))
                for item in shared_root.findall("x:si", spreadsheet_ns)
            ]
            sheet_root = ET.fromstring(archive.read("xl/worksheets/sheet2.xml"))
        xlsx_rows = []
        for row in sheet_root.findall(".//x:sheetData/x:row", spreadsheet_ns):
            values = []
            for cell in row.findall("x:c", spreadsheet_ns):
                value = cell.find("x:v", spreadsheet_ns)
                self.assertIsNotNone(value)
                raw = value.text or ""
                values.append(shared_strings[int(raw)] if cell.get("t") == "s" else raw)
            xlsx_rows.append(values)

        expected = [
            ["Label", "Detail"],
            ["Qualified | converted", "Retained"],
        ]
        self.assertEqual(html_rows, expected)
        self.assertEqual(docx_rows, expected)
        self.assertEqual(pptx_rows, expected)
        # XLSX appends its existing Notes area below native table cells. Compare
        # the table's artifact rows, not that separate prose region.
        self.assertEqual(xlsx_rows[:len(expected)], expected)


class TestORenderPPTX(_TempWorkspace):

    SAMPLE_PPTX_OFF = dedent("""\
        ---
        name: monthly-board-deck
        medium: pptx
        title: "Monthly Board Deck — {period}"
        intro: "This deck reviews monthly performance."
        sections:
          - section: weekly_sales
            heading: Weekly Sales
          - section: campaigns
            heading: Campaign Performance
        ---
        """)

    def setUp(self):
        super().setUp()
        # Replace the markdown OFF with a pptx-medium OFF
        self.off_path = os.path.join(self.tmp, "off-pptx.md")
        with open(self.off_path, "w") as f:
            f.write(self.SAMPLE_PPTX_OFF)
        result = c_instance(self.template_path, "2026-05", self.instance_dir)
        self.instance_path = result.instance_path
        # Populate both sections so the deck has content
        with open(self.instance_path) as f:
            content = f.read()
        content = content.replace(
            "<!-- Section content goes here. -->",
            "Sales were strong this month. Top product hit $42M.",
            1,
        )
        content = content.replace(
            "<!-- Section content goes here. -->",
            "- Campaign A: 1.2M impressions\n- Campaign B: 800K impressions\n- Campaign C: not yet reported",
            1,
        )
        with open(self.instance_path, "w") as f:
            f.write(content)

    def test_renders_pptx_artifact(self):
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        self.assertTrue(result.success, result.error)
        self.assertEqual(result.artifact_format, "pptx")
        self.assertTrue(result.artifact_path.endswith(".pptx"))
        self.assertTrue(os.path.isfile(result.artifact_path))

    def test_pptx_has_title_intro_and_section_slides(self):
        from pptx import Presentation
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        prs = Presentation(result.artifact_path)
        # Expect: title slide + intro slide + one slide per corpus section
        self.assertGreaterEqual(len(prs.slides), 4)

        titles = [
            s.shapes.title.text if s.shapes.title is not None else ""
            for s in prs.slides
        ]
        self.assertIn("Monthly Board Deck — 2026-05", titles[0])
        self.assertIn("Overview", titles)
        self.assertIn("Weekly Sales", titles)
        self.assertIn("Campaign Performance", titles)

    def test_pptx_bullets_render_as_separate_paragraphs(self):
        from pptx import Presentation
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        prs = Presentation(result.artifact_path)
        # Find the Campaign Performance slide
        campaign_slide = next(
            s for s in prs.slides
            if s.shapes.title is not None and "Campaign" in s.shapes.title.text
        )
        body_text_pieces = []
        for shape in campaign_slide.shapes:
            if shape.has_text_frame and shape != campaign_slide.shapes.title:
                for p in shape.text_frame.paragraphs:
                    body_text_pieces.append(p.text)
        all_body = "\n".join(body_text_pieces)
        # Each bullet line lands as its own paragraph
        self.assertIn("Campaign A", all_body)
        self.assertIn("Campaign B", all_body)
        self.assertIn("Campaign C", all_body)
        # And we have at least three non-empty paragraphs (one per bullet)
        non_empty = [p for p in body_text_pieces if p.strip()]
        self.assertGreaterEqual(len(non_empty), 3)

    def test_fenced_h2_creates_no_slide_and_following_prose_stays_in_section(self):
        from pptx import Presentation

        # Put the fenced Markdown in the OFF section intro so it reaches the
        # shared presentation/spreadsheet splitter. Corpus parsing has its own
        # heading boundary and is outside the behavior under test here.
        off_text = self.SAMPLE_PPTX_OFF.replace(
            "    heading: Weekly Sales\n",
            (
                "    heading: Weekly Sales\n"
                "    intro: |\n"
                "      Before the fence.\n"
                "      ```markdown\n"
                "      ## Fenced Heading\n"
                "      code body\n"
                "      ```\n"
                "      After the closing fence remains prose.\n"
            ),
            1,
        )
        with open(self.off_path, "w") as f:
            f.write(off_text)

        result = o_render(self.off_path, self.instance_path, self.output_dir)
        self.assertTrue(result.success, result.error)
        prs = Presentation(result.artifact_path)
        titles = [
            slide.shapes.title.text if slide.shapes.title is not None else ""
            for slide in prs.slides
        ]
        self.assertNotIn("Fenced Heading", titles)

        weekly_slide = next(
            slide for slide in prs.slides
            if slide.shapes.title is not None and slide.shapes.title.text == "Weekly Sales"
        )
        weekly_paragraphs = [
            paragraph
            for shape in weekly_slide.shapes
            if shape.has_text_frame and shape != weekly_slide.shapes.title
            for paragraph in shape.text_frame.paragraphs
        ]
        weekly_text = "\n".join(paragraph.text for paragraph in weekly_paragraphs)
        self.assertIn("## Fenced Heading", weekly_text)
        following_prose = next(
            paragraph
            for paragraph in weekly_paragraphs
            if paragraph.text == "After the closing fence remains prose."
        )
        self.assertTrue(all(run.font.name != "Courier New" for run in following_prose.runs))

    def test_unknown_medium_returns_error(self):
        # Sanity: medium=ppt (not pptx) should be rejected by the dispatcher
        bogus = os.path.join(self.tmp, "off-bogus.md")
        with open(bogus, "w") as f:
            f.write(dedent("""\
                ---
                name: bogus
                medium: ppt
                sections:
                  - section: weekly_sales
                    heading: Weekly Sales
                ---
                """))
        result = o_render(bogus, self.instance_path, self.output_dir)
        self.assertFalse(result.success)
        self.assertIn("supported:", result.error)


class TestORenderXLSX(_TempWorkspace):

    SAMPLE_XLSX_OFF = dedent("""\
        ---
        name: quarterly-financials
        medium: xlsx
        title: "Quarterly Financials — {period}"
        intro: "Q1 financial summary across product lines."
        sections:
          - section: weekly_sales
            heading: Revenue Lines
          - section: campaigns
            heading: Notes
        ---
        """)

    def setUp(self):
        super().setUp()
        # Replace the markdown OFF with an xlsx-medium OFF
        self.off_path = os.path.join(self.tmp, "off-xlsx.md")
        with open(self.off_path, "w") as f:
            f.write(self.SAMPLE_XLSX_OFF)
        result = c_instance(self.template_path, "2026-Q1", self.instance_dir)
        self.instance_path = result.instance_path
        # Populate first section with a markdown table; second with bullets
        with open(self.instance_path) as f:
            content = f.read()
        content = content.replace(
            "<!-- Section content goes here. -->",
            (
                "Quarterly revenue breakdown.\n\n"
                "| Product | Q4 2025 | Q1 2026 | Growth |\n"
                "|---------|---------|---------|--------|\n"
                "| Alpha   | $12.5M  | $14.2M  | 13.6%  |\n"
                "| Beta    | $8.0M   | $9.1M   | 13.8%  |\n"
                "| Gamma   | $5.0M   | $4.8M   | -4.0%  |\n"
            ),
            1,
        )
        content = content.replace(
            "<!-- Section content goes here. -->",
            "- Strong demand in Alpha.\n- Beta refresh launched in March.",
            1,
        )
        with open(self.instance_path, "w") as f:
            f.write(content)

    def test_renders_xlsx_artifact(self):
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        self.assertTrue(result.success, result.error)
        self.assertEqual(result.artifact_format, "xlsx")
        self.assertTrue(result.artifact_path.endswith(".xlsx"))
        self.assertTrue(os.path.isfile(result.artifact_path))

    def test_xlsx_has_overview_and_section_worksheets(self):
        # Inspect via zip + shared strings; works without openpyxl
        import zipfile
        import re as _re
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        with zipfile.ZipFile(result.artifact_path) as z:
            workbook_xml = z.read("xl/workbook.xml").decode()
        sheet_names = _re.findall(r'<sheet[^>]+name="([^"]+)"', workbook_xml)
        self.assertIn("Overview", sheet_names)
        self.assertIn("Revenue Lines", sheet_names)
        self.assertIn("Notes", sheet_names)

    def test_xlsx_table_cells_and_numeric_coercion(self):
        import zipfile
        import re as _re
        result = o_render(self.off_path, self.instance_path, self.output_dir)
        with zipfile.ZipFile(result.artifact_path) as z:
            ss = z.read("xl/sharedStrings.xml").decode()
            sheets = z.namelist()
            revenue_sheet_xml = ""
            for s in sheets:
                if s.startswith("xl/worksheets/sheet") and s.endswith(".xml"):
                    body = z.read(s).decode()
                    if "Product" in self._lookup_strings(ss, body):
                        revenue_sheet_xml = body
                        break
        # Header strings present
        all_strings = _re.findall(r"<t[^>]*>([^<]+)</t>", ss)
        for header in ["Product", "Q4 2025", "Q1 2026", "Growth"]:
            self.assertIn(header, all_strings)
        # Percentage strings were coerced to native numerics (13.6% → 0.136)
        numerics = _re.findall(r"<v>(-?\d+\.\d+|\d+)</v>", revenue_sheet_xml)
        floats = [float(n) for n in numerics if "." in n]
        self.assertTrue(any(abs(f - 0.136) < 1e-6 for f in floats),
                        f"expected 0.136 in numeric cells; got {floats}")
        self.assertTrue(any(abs(f - (-0.04)) < 1e-6 for f in floats),
                        f"expected -0.04 in numeric cells; got {floats}")
        # Currency strings ($12.5M etc.) stayed as strings (alphabetic char rejects coercion)
        self.assertIn("$12.5M", all_strings)
        self.assertIn("$14.2M", all_strings)

    def test_xlsx_long_section_heading_truncated_to_31_chars(self):
        # Excel sheet names cap at 31 chars; the renderer must not error
        long_off = os.path.join(self.tmp, "off-xlsx-long.md")
        with open(long_off, "w") as f:
            f.write(dedent("""\
                ---
                name: long-headings
                medium: xlsx
                sections:
                  - section: weekly_sales
                    heading: A Very Long Section Heading That Definitely Exceeds The Excel Limit
                ---
                """))
        result = o_render(long_off, self.instance_path, self.output_dir)
        self.assertTrue(result.success, result.error)

    def test_fenced_h2_creates_no_sheet_and_following_prose_stays_in_section(self):
        import re as _re
        import zipfile

        off_text = self.SAMPLE_XLSX_OFF.replace(
            "    heading: Revenue Lines\n",
            (
                "    heading: Revenue Lines\n"
                "    intro: |\n"
                "      Before the fence.\n"
                "      ```markdown\n"
                "      ## Fenced Heading\n"
                "      code body\n"
                "      ```\n"
                "      After the closing fence remains prose.\n"
            ),
            1,
        )
        with open(self.off_path, "w") as f:
            f.write(off_text)

        result = o_render(self.off_path, self.instance_path, self.output_dir)
        self.assertTrue(result.success, result.error)
        with zipfile.ZipFile(result.artifact_path) as archive:
            workbook_xml = archive.read("xl/workbook.xml").decode()
            shared_strings = archive.read("xl/sharedStrings.xml").decode()
            revenue_sheet = archive.read("xl/worksheets/sheet2.xml").decode()

        sheet_names = _re.findall(r'<sheet[^>]+name="([^"]+)"', workbook_xml)
        self.assertNotIn("Fenced Heading", sheet_names)
        revenue_strings = self._lookup_strings(shared_strings, revenue_sheet)
        self.assertIn("## Fenced Heading", revenue_strings)
        self.assertIn("After the closing fence remains prose.", revenue_strings)

    def test_unsupported_xlsx_alias_still_rejected(self):
        # `xls` (legacy) is not in the supported list — must reject cleanly
        bogus = os.path.join(self.tmp, "off-xls.md")
        with open(bogus, "w") as f:
            f.write(dedent("""\
                ---
                name: bogus-xls
                medium: xls
                sections:
                  - section: weekly_sales
                    heading: Weekly Sales
                ---
                """))
        result = o_render(bogus, self.instance_path, self.output_dir)
        self.assertFalse(result.success)
        self.assertIn("supported:", result.error)

    @staticmethod
    def _lookup_strings(shared_strings_xml: str, sheet_xml: str) -> set:
        """Resolve numeric-indexed string references in a sheet against the
        sharedStrings table. Returns the set of distinct string values that
        appear as cells in the sheet."""
        import re as _re
        all_strings = _re.findall(r"<t[^>]*>([^<]+)</t>", shared_strings_xml)
        # Sheet cells of type "s" carry the index into shared strings
        indices = _re.findall(r't="s"[^>]*><v>(\d+)</v>', sheet_xml)
        return {all_strings[int(i)] for i in indices if int(i) < len(all_strings)}


class TestOFFSpecParsing(unittest.TestCase):

    def test_parses_frontmatter_fields(self):
        with tempfile.NamedTemporaryFile("w", suffix=".md", delete=False) as f:
            f.write(SAMPLE_OFF)
            path = f.name
        try:
            spec = parse_off_spec(path)
            self.assertEqual(spec.name, "monthly-board-memo")
            self.assertEqual(spec.medium, "markdown")
            self.assertEqual(len(spec.sections), 2)
            self.assertEqual(spec.sections[0].corpus_section_id, "weekly_sales")
        finally:
            os.unlink(path)

    def test_malformed_or_structureless_spec_is_explicitly_invalid(self):
        cases = {
            "malformed-yaml": "---\nsections: [\n---\n",
            "missing-sections": "---\nname: empty\n---\n# No section markers\n",
        }
        for label, content in cases.items():
            with self.subTest(case=label):
                with tempfile.NamedTemporaryFile("w", suffix=".md", delete=False) as f:
                    f.write(content)
                    path = f.name
                try:
                    spec = parse_off_spec(path)
                    self.assertFalse(spec.is_valid)
                    self.assertTrue(spec.parse_error)
                finally:
                    os.unlink(path)

    def test_body_section_marker_style_remains_valid(self):
        with tempfile.NamedTemporaryFile("w", suffix=".md", delete=False) as f:
            f.write("## Section: SalesQ2\n")
            path = f.name
        try:
            spec = parse_off_spec(path)
            self.assertTrue(spec.is_valid, spec.parse_error)
            self.assertEqual([section.corpus_section_id for section in spec.sections], ["SalesQ2"])
        finally:
            os.unlink(path)


class TestEndToEndShape4(_TempWorkspace):
    """One full pass: template → instance → populated → validated → rendered."""

    def test_full_pipeline_emits_all_events(self):
        # 1. C-Instance
        ci = c_instance(self.template_path, "2026-05", self.instance_dir,
                        workflow_id="smoke", project_nexus="smoke_test")
        self.assertTrue(ci.success)

        # 2. Populate
        with open(ci.instance_path) as f:
            content = f.read()
        content = content.replace(
            "<!-- Section content goes here. -->",
            "Real data.",
        )
        with open(ci.instance_path, "w") as f:
            f.write(content)

        # 3. C-Validate
        cv = c_validate(ci.instance_path, self.template_path,
                        workflow_id="smoke", project_nexus="smoke_test")
        self.assertEqual(cv.overall_status, "PASS")

        # 4. O-Render
        orr = o_render(self.off_path, ci.instance_path, self.output_dir,
                       workflow_id="smoke", project_nexus="smoke_test")
        self.assertTrue(orr.success)

        # 5. Events
        for expected in ("CorpusInstanceCreated", "CorpusValidated", "OFFRendered"):
            self.assertIn(expected, self.events, f"Missing event: {expected}")

    def test_mixed_case_section_id_survives_instance_validation_and_render(self):
        with open(self.template_path, "w") as f:
            f.write(dedent("""\
                ---
                type: corpus_template
                template_version: 1.0
                ---

                # Quarterly Corpus Template

                ## Sections

                ### Section SalesQ2 — Quarterly Sales

                *Source PFF:* `pff-quarterly-sales`
                *Missing-data behavior:* hold-and-warn
                """))
        with open(self.off_path, "w") as f:
            f.write(dedent("""\
                ---
                name: quarterly-memo
                medium: markdown
                sections:
                  - section: SalesQ2
                    heading: Quarterly Sales
                ---
                """))

        created = c_instance(self.template_path, "2026-Q2", self.instance_dir)
        self.assertTrue(created.success, created.error)
        self.assertEqual(created.sections_created, ["SalesQ2"])

        with open(created.instance_path) as f:
            content = f.read().replace(
                "<!-- Section content goes here. -->",
                "Mixed-case identity reached the corpus body.",
            )
        with open(created.instance_path, "w") as f:
            f.write(content)

        parsed = parse_corpus_file(created.instance_path)
        self.assertEqual([s.section_id for s in parsed.sections], ["SalesQ2"])

        validated = c_validate(created.instance_path, self.template_path)
        self.assertTrue(validated.success, validated.error)
        self.assertEqual([s.section_id for s in validated.sections], ["SalesQ2"])

        rendered = o_render(self.off_path, created.instance_path, self.output_dir)
        self.assertTrue(rendered.success, rendered.error)
        self.assertEqual(rendered.sections_read, ["SalesQ2"])
        with open(rendered.artifact_path) as f:
            self.assertIn("Mixed-case identity reached the corpus body.", f.read())


if __name__ == "__main__":
    unittest.main()
