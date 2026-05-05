"""Output runtime — O-Render markdown renderer.

Implements the mechanical OFF mode that takes a corpus instance + a bespoke
OFF spec and renders a final artifact in markdown. The model-driven OFF
modes (O-Design, O-Modify, O-Audit) are not implemented here.

A bespoke OFF spec is itself a markdown file with YAML frontmatter that
declares: which sections to read from the corpus, in what order, with what
heading structure, and how to present missing-data sections. The MVP
renders to markdown only; future formats (Word, PowerPoint, etc.) can be
added as additional render functions consuming the same OFF spec.

Per Reference — Meta-Layer Architecture §5 (E10 OFFRendered).
"""
from __future__ import annotations

import os
import re
from dataclasses import dataclass, field
from datetime import datetime, timezone
from typing import Optional

try:
    import yaml
except ImportError:
    yaml = None

from corpus_parser import parse_corpus_file, ParsedCorpus, CorpusSection


# ---------- Result types ----------

@dataclass
class RenderResult:
    success: bool
    artifact_path: str = ""
    artifact_format: str = "markdown"
    sections_read: list = field(default_factory=list)
    error: str = ""


@dataclass
class OFFSpec:
    name: str = ""
    medium: str = "markdown"  # only "markdown" implemented in MVP
    title_template: str = ""
    intro_text: str = ""
    sections: list = field(default_factory=list)  # list[OFFSectionSpec]
    missing_section_behavior: str = "note"  # "note" / "skip" / "fail"
    raw_frontmatter: dict = field(default_factory=dict)


@dataclass
class OFFSectionSpec:
    corpus_section_id: str
    heading: str = ""
    heading_level: int = 2
    intro: str = ""


# ---------- O-Render ----------

def o_render(
    off_spec_path: str,
    corpus_instance_path: str,
    output_dir: str,
    workflow_id: Optional[str] = None,
    project_nexus: Optional[str] = None,
) -> RenderResult:
    """Render a corpus instance through a bespoke OFF spec to produce an artifact.

    Supported media (declared in the OFF spec's frontmatter `medium:` field):
      - "markdown" / "md" — plain markdown
      - "html" — HTML wrapping the markdown (simple)
      - "docx" / "word" — Microsoft Word document (requires python-docx)

    Args:
        off_spec_path: path to the bespoke OFF markdown file
        corpus_instance_path: path to the corpus instance to render
        output_dir: directory where the rendered artifact is written
        workflow_id: optional workflow id for the event payload
        project_nexus: optional project nexus for the event payload

    Returns: RenderResult with artifact_path on success.
    """
    if not os.path.isfile(off_spec_path):
        return RenderResult(success=False, error=f"OFF spec not found: {off_spec_path}")
    if not os.path.isfile(corpus_instance_path):
        return RenderResult(success=False, error=f"Corpus instance not found: {corpus_instance_path}")

    try:
        off_spec = parse_off_spec(off_spec_path)
    except Exception as e:
        return RenderResult(success=False, error=f"Failed to parse OFF spec: {e}")

    medium = off_spec.medium.lower()
    if medium not in ("markdown", "md", "html", "docx", "word", "pptx", "powerpoint", "xlsx", "excel"):
        return RenderResult(
            success=False,
            error=f"OFF spec declares medium={medium!r}; supported: markdown, html, docx, pptx, xlsx",
        )

    try:
        corpus = parse_corpus_file(corpus_instance_path)
    except Exception as e:
        return RenderResult(success=False, error=f"Failed to parse corpus: {e}")

    section_map = {s.section_id: s for s in corpus.sections}

    # Build the rendered output (always start with markdown lines; downstream
    # converters consume that intermediate form)
    output_lines, sections_read, missing_sections = _render_markdown(off_spec, corpus, section_map)

    if missing_sections and off_spec.missing_section_behavior == "fail":
        return RenderResult(
            success=False,
            error=f"Required sections missing from corpus: {missing_sections}",
        )

    os.makedirs(output_dir, exist_ok=True)
    artifact_filename = _artifact_filename(off_spec, corpus, off_spec_path, medium)
    artifact_path = os.path.join(output_dir, artifact_filename)

    markdown_text = "\n".join(output_lines)

    try:
        if medium in ("markdown", "md"):
            with open(artifact_path, "w", encoding="utf-8") as f:
                f.write(markdown_text)
            artifact_format = "markdown"
        elif medium == "html":
            html_text = _markdown_to_html(markdown_text, off_spec, corpus)
            with open(artifact_path, "w", encoding="utf-8") as f:
                f.write(html_text)
            artifact_format = "html"
        elif medium in ("docx", "word"):
            artifact_format = "docx"
            err = _markdown_to_docx(markdown_text, artifact_path, off_spec, corpus)
            if err:
                return RenderResult(success=False, error=err)
        elif medium in ("pptx", "powerpoint"):
            artifact_format = "pptx"
            err = _markdown_to_pptx(markdown_text, artifact_path, off_spec, corpus)
            if err:
                return RenderResult(success=False, error=err)
        elif medium in ("xlsx", "excel"):
            artifact_format = "xlsx"
            err = _markdown_to_xlsx(markdown_text, artifact_path, off_spec, corpus)
            if err:
                return RenderResult(success=False, error=err)
        else:
            # Already validated above — defensive
            return RenderResult(success=False, error=f"Unsupported medium: {medium}")
    except OSError as e:
        return RenderResult(success=False, error=f"Failed to write artifact: {e}")

    # Emit OFFRendered event
    try:
        from oversight_events import emit
        emit({
            "event_type": "OFFRendered",
            "workflow_id": workflow_id or "",
            "project_nexus": project_nexus or "",
            "corpus_instance_path": corpus_instance_path,
            "off_framework_id": off_spec.name or os.path.basename(off_spec_path),
            "sections_read": sections_read,
            "artifact_path": artifact_path,
            "artifact_format": artifact_format,
        })
    except Exception:
        pass

    return RenderResult(
        success=True,
        artifact_path=artifact_path,
        artifact_format=artifact_format,
        sections_read=sections_read,
    )


# ---------- OFF spec parsing ----------

def parse_off_spec(path: str) -> OFFSpec:
    """Parse a bespoke OFF markdown file into an OFFSpec."""
    with open(path, encoding="utf-8") as f:
        content = f.read()

    fm, body = _split_frontmatter(content)
    spec = OFFSpec()

    if fm and yaml is not None:
        try:
            data = yaml.safe_load(fm) or {}
        except yaml.YAMLError:
            data = {}
        if isinstance(data, dict):
            spec.raw_frontmatter = data
            spec.name = str(data.get("name", os.path.splitext(os.path.basename(path))[0]))
            spec.medium = str(data.get("medium", "markdown")).lower()
            spec.title_template = str(data.get("title", ""))
            spec.intro_text = str(data.get("intro", ""))
            spec.missing_section_behavior = str(data.get("missing_section_behavior", "note"))

            # Sections list may live in frontmatter
            for s in data.get("sections", []) or []:
                if not isinstance(s, dict):
                    continue
                spec.sections.append(OFFSectionSpec(
                    corpus_section_id=str(s.get("section", "")),
                    heading=str(s.get("heading", "")),
                    heading_level=int(s.get("heading_level", 2) or 2),
                    intro=str(s.get("intro", "")),
                ))

    if not spec.name:
        spec.name = os.path.splitext(os.path.basename(path))[0]

    # If no sections in frontmatter, fall back to scanning the body for
    # "## Section: <id>" markers (an alternate authoring style)
    if not spec.sections:
        for m in re.finditer(r"^##\s+Section:\s*(.+?)\s*$", body, re.MULTILINE):
            spec.sections.append(OFFSectionSpec(
                corpus_section_id=m.group(1).strip(),
                heading=m.group(1).strip().replace("_", " ").title(),
                heading_level=2,
            ))

    return spec


def _split_frontmatter(content: str) -> tuple[str, str]:
    pat = re.compile(r"^---\s*\n(.*?)\n---\s*\n", re.DOTALL)
    m = pat.match(content)
    if not m:
        return ("", content)
    return (m.group(1), content[m.end():])


# ---------- Markdown rendering ----------

def _render_markdown(
    off_spec: OFFSpec,
    corpus: ParsedCorpus,
    section_map: dict[str, CorpusSection],
) -> tuple[list[str], list[str], list[str]]:
    """Build the rendered markdown lines. Returns (lines, sections_read, missing)."""
    lines: list[str] = []
    sections_read: list[str] = []
    missing_sections: list[str] = []

    # Title
    title = off_spec.title_template
    if not title:
        base = re.sub(r"^Corpus Template[-\s—:]+", "", corpus.title or "", flags=re.IGNORECASE).strip()
        title = base or "Rendered Output"
    title = _interpolate(title, {"period": corpus.instance_period or "", "corpus_title": corpus.title or ""})
    lines.append(f"# {title}")
    lines.append("")
    lines.append(f"*Rendered: {_now_iso()}*")
    if corpus.instance_period:
        lines.append(f"*Period: {corpus.instance_period}*")
    if off_spec.name:
        lines.append(f"*Generated by OFF:* `{off_spec.name}`")
    lines.append("")

    if off_spec.intro_text:
        lines.append(off_spec.intro_text)
        lines.append("")

    # Sections
    for section_spec in off_spec.sections:
        corpus_section = section_map.get(section_spec.corpus_section_id)
        heading = section_spec.heading or section_spec.corpus_section_id
        level_marker = "#" * max(1, min(section_spec.heading_level, 6))
        lines.append(f"{level_marker} {heading}")
        lines.append("")

        if section_spec.intro:
            lines.append(section_spec.intro)
            lines.append("")

        if corpus_section is None:
            missing_sections.append(section_spec.corpus_section_id)
            if off_spec.missing_section_behavior == "skip":
                lines.append("*(section not present in corpus — skipped)*")
            else:
                lines.append("*(section not present in corpus)*")
            lines.append("")
            continue

        body = (corpus_section.raw_body or "").strip()
        # Strip the section's own meta lines (matched by *Source PFF:* etc. patterns)
        body = re.sub(r"^\*[^*]+:\*\s.*$", "", body, flags=re.MULTILINE)
        body = re.sub(r"<!--.*?-->", "", body, flags=re.DOTALL).strip()

        if not body:
            lines.append("*(section is empty)*")
        else:
            lines.append(body)
        lines.append("")
        sections_read.append(section_spec.corpus_section_id)

    return (lines, sections_read, missing_sections)


def _interpolate(template: str, values: dict[str, str]) -> str:
    """Simple {key} interpolation."""
    out = template
    for k, v in values.items():
        out = out.replace("{" + k + "}", str(v))
    return out


def _artifact_filename(off_spec: OFFSpec, corpus: ParsedCorpus, off_spec_path: str, medium: str = "markdown") -> str:
    base = re.sub(r"[^a-z0-9]+", "-", (off_spec.name or "artifact").lower()).strip("-")
    period = re.sub(r"[^A-Za-z0-9_-]+", "_", corpus.instance_period or "")
    ext = {
        "markdown": "md", "md": "md",
        "html": "html",
        "docx": "docx", "word": "docx",
        "pptx": "pptx", "powerpoint": "pptx",
        "xlsx": "xlsx", "excel": "xlsx",
    }.get(medium, "md")
    if period:
        return f"{base}-{period}.{ext}"
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    return f"{base}-{stamp}.{ext}"


# ---------- HTML rendering (simple) ----------

def _markdown_to_html(md_text: str, off_spec: OFFSpec, corpus: ParsedCorpus) -> str:
    """Tiny markdown-to-HTML converter. Handles headings, paragraphs, lists,
    tables, code blocks. No external dependency. Output is simple but valid.
    """
    title = re.sub(r"^#\s+", "", md_text.split("\n")[0]) if md_text.strip() else "Rendered Output"
    body = _md_to_html_body(md_text)
    return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{_html_escape(title)}</title>
<style>
body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; max-width: 760px; margin: 2em auto; padding: 0 1em; color: #222; }}
h1, h2, h3 {{ color: #111; }}
table {{ border-collapse: collapse; }}
th, td {{ border: 1px solid #ccc; padding: 0.4em 0.7em; }}
th {{ background: #f0f0f0; }}
pre {{ background: #f6f6f6; padding: 0.8em; overflow-x: auto; border-radius: 4px; }}
em {{ color: #555; }}
</style>
</head>
<body>
{body}
</body>
</html>
"""


def _md_to_html_body(md: str) -> str:
    lines = md.split("\n")
    out: list[str] = []
    in_table = False
    table_rows: list[list[str]] = []
    in_code = False
    paragraph: list[str] = []

    def flush_paragraph():
        if paragraph:
            text = " ".join(paragraph).strip()
            if text:
                # Bold and italic
                text = _md_inline_to_html(text)
                out.append(f"<p>{text}</p>")
            paragraph.clear()

    def flush_table():
        nonlocal in_table
        if not table_rows:
            return
        out.append("<table>")
        # First row as header, second row is the divider, rest are data
        for i, row in enumerate(table_rows):
            if i == 1 and all(re.match(r"^[-:\s]+$", c.strip()) for c in row):
                continue
            tag = "th" if i == 0 else "td"
            cells = "".join(f"<{tag}>{_md_inline_to_html(_html_escape(c.strip()))}</{tag}>" for c in row)
            out.append(f"<tr>{cells}</tr>")
        out.append("</table>")
        table_rows.clear()
        in_table = False

    for line in lines:
        if line.startswith("```"):
            flush_paragraph()
            flush_table()
            if not in_code:
                in_code = True
                out.append("<pre><code>")
            else:
                in_code = False
                out.append("</code></pre>")
            continue
        if in_code:
            out.append(_html_escape(line))
            continue

        # Tables — pipe-separated
        if "|" in line and line.strip().startswith("|"):
            flush_paragraph()
            in_table = True
            cells = [c for c in line.strip().strip("|").split("|")]
            table_rows.append(cells)
            continue
        if in_table:
            flush_table()

        # Headings
        m = re.match(r"^(#{1,6})\s+(.+)$", line)
        if m:
            flush_paragraph()
            level = len(m.group(1))
            out.append(f"<h{level}>{_md_inline_to_html(_html_escape(m.group(2)))}</h{level}>")
            continue

        # Bullets
        m = re.match(r"^\s*[-*]\s+(.+)$", line)
        if m:
            flush_paragraph()
            if not out or not out[-1].endswith("</li>"):
                out.append("<ul>")
            out.append(f"<li>{_md_inline_to_html(_html_escape(m.group(1)))}</li>")
            continue
        elif out and out[-1].startswith("<li>"):
            out.append("</ul>")

        if line.strip() == "":
            flush_paragraph()
            continue

        paragraph.append(line)

    flush_paragraph()
    flush_table()
    if out and out[-1].startswith("<li>"):
        out.append("</ul>")
    if in_code:
        out.append("</code></pre>")

    return "\n".join(out)


def _md_inline_to_html(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", s)
    s = re.sub(r"\*(.+?)\*", r"<em>\1</em>", s)
    s = re.sub(r"`(.+?)`", r"<code>\1</code>", s)
    return s


def _html_escape(s: str) -> str:
    return (s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


# ---------- DOCX rendering (via python-docx) ----------

def _markdown_to_docx(md_text: str, output_path: str, off_spec: OFFSpec, corpus: ParsedCorpus) -> str:
    """Convert markdown to a .docx file. Returns "" on success, error string on failure."""
    try:
        from docx import Document  # type: ignore
        from docx.shared import Pt  # type: ignore
    except ImportError:
        return (
            "DOCX rendering requires python-docx. Install with: "
            "/opt/homebrew/bin/python3 -m pip install python-docx"
        )

    doc = Document()

    lines = md_text.split("\n")
    in_code = False
    code_buffer: list[str] = []
    table_rows: list[list[str]] = []
    in_table = False
    paragraph_buffer: list[str] = []

    def flush_paragraph():
        if paragraph_buffer:
            text = " ".join(paragraph_buffer).strip()
            if text:
                _docx_add_paragraph(doc, text)
            paragraph_buffer.clear()

    def flush_table():
        nonlocal in_table
        if not table_rows:
            return
        # Skip alignment row
        rows = [r for r in table_rows if not all(re.match(r"^[-:\s]+$", c.strip()) for c in r)]
        if rows:
            tbl = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
            tbl.style = "Light Grid Accent 1"
            for i, r in enumerate(rows):
                for j, c in enumerate(r):
                    if j < len(tbl.rows[i].cells):
                        tbl.rows[i].cells[j].text = c.strip()
        table_rows.clear()
        in_table = False

    for line in lines:
        if line.startswith("```"):
            flush_paragraph()
            flush_table()
            if in_code:
                # End of code block
                p = doc.add_paragraph()
                run = p.add_run("\n".join(code_buffer))
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                code_buffer.clear()
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_buffer.append(line)
            continue

        # Tables
        if "|" in line and line.strip().startswith("|"):
            flush_paragraph()
            in_table = True
            table_rows.append([c for c in line.strip().strip("|").split("|")])
            continue
        if in_table:
            flush_table()

        m = re.match(r"^(#{1,6})\s+(.+)$", line)
        if m:
            flush_paragraph()
            level = min(len(m.group(1)), 4)
            doc.add_heading(_strip_inline_md(m.group(2)), level=level)
            continue

        m = re.match(r"^\s*[-*]\s+(.+)$", line)
        if m:
            flush_paragraph()
            doc.add_paragraph(_strip_inline_md(m.group(1)), style="List Bullet")
            continue

        if line.strip() == "":
            flush_paragraph()
            continue

        paragraph_buffer.append(line)

    flush_paragraph()
    flush_table()

    try:
        doc.save(output_path)
    except Exception as e:
        return f"Failed to save DOCX: {e}"
    return ""


def _docx_add_paragraph(doc, text: str):
    """Add a paragraph with simple inline-markdown handling (bold, italic, code)."""
    text = text.strip()
    p = doc.add_paragraph()
    # Walk the text and emit runs with appropriate formatting
    pos = 0
    pattern = re.compile(r"(\*\*.+?\*\*|\*[^*]+?\*|`[^`]+?`)")
    for m in pattern.finditer(text):
        if m.start() > pos:
            p.add_run(text[pos:m.start()])
        chunk = m.group(0)
        if chunk.startswith("**"):
            run = p.add_run(chunk[2:-2])
            run.bold = True
        elif chunk.startswith("*"):
            run = p.add_run(chunk[1:-1])
            run.italic = True
        elif chunk.startswith("`"):
            run = p.add_run(chunk[1:-1])
            run.font.name = "Courier New"
        pos = m.end()
    if pos < len(text):
        p.add_run(text[pos:])


def _strip_inline_md(text: str) -> str:
    """Strip simple inline markdown (used for headings and list items in DOCX)."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text


# ---------- PPTX rendering (via python-pptx) ----------

# Slide-mapping conventions:
#   - The first H1 in the rendered markdown becomes the title slide. Italicized
#     metadata lines that follow (Rendered: ..., Period: ..., Generated by OFF: ...)
#     stack into the subtitle placeholder.
#   - The OFF spec's intro_text, when present, becomes a dedicated intro slide.
#   - Each H2 boundary opens a new content slide. The H2 text is the slide title;
#     the body until the next H2 fills the content placeholder.
#   - Inside a content slide: H3 lines render as bold lead-in paragraphs; bullets
#     render as bullet-level paragraphs; ordinary lines render as paragraphs;
#     pipe tables render as native PPTX tables added below the text frame; fenced
#     code blocks render as a monospaced text frame.
#
# This is a workable first cut, not a presentation-design framework. The goal is
# producing a usable PPTX from a corpus instance + bespoke OFF without further
# manual cleanup; sophisticated layout decisions belong to the bespoke OFF spec
# and would require additional render-layer hints (deferred).

def _markdown_to_pptx(md_text: str, output_path: str, off_spec: OFFSpec, corpus: ParsedCorpus) -> str:
    """Convert markdown to a .pptx file. Returns "" on success, error string on failure."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except ImportError:
        return (
            "PPTX rendering requires python-pptx. Install with: "
            "/opt/homebrew/bin/python3 -m pip install python-pptx"
        )

    prs = Presentation()

    blocks = _split_md_into_render_blocks(md_text)
    if not blocks:
        return "No renderable content found in markdown."

    title_block = blocks[0]
    _add_title_slide(prs, title_block, off_spec, corpus)

    # Optional intro slide if the OFF spec declared intro_text and the title
    # block didn't already absorb it
    if off_spec.intro_text and off_spec.intro_text.strip() not in title_block.subtitle_lines:
        _add_intro_slide(prs, off_spec.intro_text)

    for block in blocks[1:]:
        _add_content_slide(prs, block)

    try:
        prs.save(output_path)
    except Exception as e:
        return f"Failed to save PPTX: {e}"
    return ""


@dataclass
class _RenderBlock:
    """A single slide's worth of content, parsed from the markdown intermediate form."""
    title: str = ""
    subtitle_lines: list = field(default_factory=list)  # used only on the title block
    body_lines: list = field(default_factory=list)      # raw body lines for content slides


def _split_md_into_render_blocks(md_text: str) -> list:
    """Walk the rendered markdown and split into slide-shaped blocks.

    The first block holds the H1 title plus the italicized metadata lines that
    immediately follow (which become the title-slide subtitle). Subsequent
    blocks each start at an H2 boundary.
    """
    blocks: list[_RenderBlock] = []
    current: _RenderBlock = _RenderBlock()
    seen_h1 = False
    in_subtitle_zone = False

    lines = md_text.split("\n")
    for raw in lines:
        line = raw.rstrip()

        h1 = re.match(r"^#\s+(.+)$", line)
        if h1 and not seen_h1:
            current.title = _strip_inline_md(h1.group(1))
            seen_h1 = True
            in_subtitle_zone = True
            continue

        h2 = re.match(r"^##\s+(.+)$", line)
        if h2:
            blocks.append(current)
            current = _RenderBlock(title=_strip_inline_md(h2.group(1)))
            in_subtitle_zone = False
            continue

        if in_subtitle_zone:
            # Italicized metadata lines feed the title subtitle until we hit
            # the first non-meta, non-blank line — at which point we leave the
            # subtitle zone and accumulate body lines instead.
            stripped = line.strip()
            if not stripped:
                continue
            meta = re.match(r"^\*([^*]+)\*$", stripped)
            if meta:
                current.subtitle_lines.append(meta.group(1).strip())
                continue
            # First substantive non-meta line — treat as body of the title block
            in_subtitle_zone = False

        current.body_lines.append(line)

    blocks.append(current)
    # Drop a leading empty block (occurs when the markdown starts without an H1)
    if blocks and not blocks[0].title and not blocks[0].body_lines and not blocks[0].subtitle_lines:
        blocks = blocks[1:]
    return blocks


def _add_title_slide(prs, block: _RenderBlock, off_spec: OFFSpec, corpus: ParsedCorpus):
    """First slide: title placeholder + subtitle placeholder."""
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = block.title or off_spec.name or "Rendered Output"

    # Subtitle placeholder — index 1 in the standard Title Slide layout
    subtitle_shape = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            subtitle_shape = ph
            break
    if subtitle_shape is None:
        return
    subtitle_text = "\n".join(block.subtitle_lines) if block.subtitle_lines else ""
    if not subtitle_text and corpus.instance_period:
        subtitle_text = f"Period: {corpus.instance_period}"
    subtitle_shape.text = subtitle_text


def _add_intro_slide(prs, intro_text: str):
    """Optional intro slide carrying the OFF spec's intro_text."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Overview"
    body_shape = _find_content_placeholder(slide)
    if body_shape is None:
        return
    tf = body_shape.text_frame
    tf.text = ""
    for i, para in enumerate(intro_text.strip().split("\n\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_paragraph_text(p, _strip_inline_md(para.strip()))


def _add_content_slide(prs, block: _RenderBlock):
    """One slide per H2 section. Bullets stay bullets; tables get appended as shapes."""
    from pptx.util import Inches, Pt  # type: ignore

    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = block.title or "(untitled section)"

    body_shape = _find_content_placeholder(slide)
    if body_shape is None:
        return

    tf = body_shape.text_frame
    tf.word_wrap = True
    # The body placeholder starts with one empty paragraph; we replace its content
    # rather than adding a new paragraph for the first line.
    first_paragraph_used = False

    table_buffers: list[list[list[str]]] = []
    current_table: list[list[str]] = []
    in_code = False
    code_buffer: list[str] = []

    def flush_code():
        nonlocal in_code, code_buffer, first_paragraph_used
        if not code_buffer:
            return
        text = "\n".join(code_buffer)
        p = tf.paragraphs[0] if not first_paragraph_used else tf.add_paragraph()
        first_paragraph_used = True
        run = p.add_run()
        run.text = text
        run.font.name = "Courier New"
        run.font.size = Pt(11)
        code_buffer = []

    def flush_table():
        nonlocal current_table
        if current_table:
            # Drop any markdown alignment row
            rows = [
                r for r in current_table
                if not all(re.match(r"^[-:\s]+$", c.strip()) for c in r)
            ]
            if rows:
                table_buffers.append(rows)
            current_table = []

    for line in block.body_lines:
        if line.startswith("```"):
            flush_table()
            if in_code:
                flush_code()
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_buffer.append(line)
            continue

        if "|" in line and line.strip().startswith("|"):
            current_table.append([c for c in line.strip().strip("|").split("|")])
            continue
        if current_table:
            flush_table()

        h3 = re.match(r"^###\s+(.+)$", line)
        if h3:
            p = tf.paragraphs[0] if not first_paragraph_used else tf.add_paragraph()
            first_paragraph_used = True
            run = p.add_run()
            run.text = _strip_inline_md(h3.group(1))
            run.font.bold = True
            continue

        bullet = re.match(r"^\s*[-*]\s+(.+)$", line)
        if bullet:
            p = tf.paragraphs[0] if not first_paragraph_used else tf.add_paragraph()
            first_paragraph_used = True
            p.level = 0
            _set_paragraph_text(p, _strip_inline_md(bullet.group(1)))
            continue

        if line.strip() == "":
            # Skip blank lines — we don't emit empty paragraphs into slides
            continue

        p = tf.paragraphs[0] if not first_paragraph_used else tf.add_paragraph()
        first_paragraph_used = True
        _set_paragraph_text(p, _strip_inline_md(line))

    if in_code:
        flush_code()
    flush_table()

    # Append any tables as native PPTX tables, stacked beneath the body shape.
    # Each table gets its own row of shapes so we don't fight the placeholder.
    if table_buffers:
        slide_w = prs.slide_width
        # Anchor below the body placeholder (rough heuristic — exact bounds
        # would require reading the layout placeholder's height)
        top = Inches(5.0)
        for tbl_rows in table_buffers:
            cols = max(len(r) for r in tbl_rows)
            rows = len(tbl_rows)
            left = Inches(0.5)
            width = slide_w - Inches(1.0)
            height = Inches(0.4 * rows)
            shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            for i, row in enumerate(tbl_rows):
                for j, cell_text in enumerate(row):
                    if j < cols:
                        shape.table.rows[i].cells[j].text = _strip_inline_md(cell_text.strip())
            top = top + height + Inches(0.3)


def _find_content_placeholder(slide):
    """Return the slide's content placeholder (idx 1 in the Title+Content layout)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    # Fallback: any non-title shape with a text frame
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            return shape
    return None


def _set_paragraph_text(p, text: str):
    """Set a paragraph's text, replacing whatever was there. Clears existing runs."""
    # python-pptx's Paragraph has no clear API; we mutate via the underlying XML
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = text


# ---------- XLSX rendering (via xlsxwriter) ----------

# Worksheet-mapping conventions:
#   - The first H1 + italicized metadata become an "Overview" worksheet at index 0
#     with title in row 1 and metadata as italic notes below.
#   - Each H2 boundary becomes a new worksheet named after the heading
#     (sanitized to ≤31 chars, no Excel-illegal characters, deduped).
#   - Within a section worksheet: any markdown pipe tables render as native cells
#     (header row formatted with bold + light-blue background, data rows bordered).
#     Non-table prose lines render below the tables in column A as italic notes.
#     Bullets render as italic lines too — this medium is for tabular data, so
#     prose is preserved as commentary rather than re-formatted.
#
# Workable first cut, not a financial-modeling toolkit. Sophisticated layouts
# (named ranges, formulas, charts) belong in the bespoke OFF spec's render layer
# and are deferred.

def _markdown_to_xlsx(md_text: str, output_path: str, off_spec: OFFSpec, corpus: ParsedCorpus) -> str:
    """Convert markdown to an .xlsx file. Returns "" on success, error string on failure."""
    try:
        import xlsxwriter  # type: ignore
    except ImportError:
        return (
            "XLSX rendering requires xlsxwriter. Install with: "
            "/opt/homebrew/bin/python3 -m pip install xlsxwriter"
        )

    blocks = _split_md_into_render_blocks(md_text)
    if not blocks:
        return "No renderable content found in markdown."

    try:
        wb = xlsxwriter.Workbook(output_path)
    except Exception as e:
        return f"Failed to open XLSX for writing: {e}"

    title_fmt = wb.add_format({"bold": True, "font_size": 14})
    meta_fmt = wb.add_format({"italic": True, "font_color": "#555555"})
    header_fmt = wb.add_format({"bold": True, "bg_color": "#D9E1F2", "border": 1})
    cell_fmt = wb.add_format({"border": 1})
    note_fmt = wb.add_format({"italic": True, "font_color": "#555555"})

    used_sheet_names: set = set()

    title_block = blocks[0]
    _add_xlsx_overview_sheet(
        wb, title_block, off_spec, corpus,
        title_fmt, meta_fmt, used_sheet_names,
    )

    for block in blocks[1:]:
        _add_xlsx_section_sheet(
            wb, block, header_fmt, cell_fmt, note_fmt, used_sheet_names,
        )

    try:
        wb.close()
    except Exception as e:
        return f"Failed to save XLSX: {e}"
    return ""


def _xlsx_safe_sheet_name(name: str, used: set) -> str:
    """Sanitize a string into a valid Excel sheet name.

    Excel rules: ≤31 chars, no `: \\ / ? * [ ]`. Names must be unique within
    a workbook. If a collision occurs, append " (N)" with truncation.
    """
    invalid = set(":\\/?*[]")
    safe = "".join(c if c not in invalid else "_" for c in name).strip()
    if not safe:
        safe = "Sheet"
    safe = safe[:31]
    base = safe
    n = 2
    while safe.lower() in {s.lower() for s in used}:
        suffix = f" ({n})"
        safe = base[: 31 - len(suffix)] + suffix
        n += 1
    used.add(safe)
    return safe


def _add_xlsx_overview_sheet(
    wb, block: _RenderBlock, off_spec: OFFSpec, corpus: ParsedCorpus,
    title_fmt, meta_fmt, used: set,
):
    """First worksheet — title + metadata + intro_text."""
    name = _xlsx_safe_sheet_name("Overview", used)
    ws = wb.add_worksheet(name)
    ws.set_column(0, 0, 80)

    title = block.title or off_spec.name or "Rendered Output"
    ws.write_string(0, 0, title, title_fmt)

    row = 1
    for line in block.subtitle_lines:
        ws.write_string(row, 0, line, meta_fmt)
        row += 1

    if corpus.instance_period and not any(
        "Period" in s for s in block.subtitle_lines
    ):
        ws.write_string(row, 0, f"Period: {corpus.instance_period}", meta_fmt)
        row += 1

    if off_spec.intro_text:
        row += 1  # blank row
        for para in off_spec.intro_text.strip().split("\n\n"):
            ws.write_string(row, 0, _strip_inline_md(para.strip()))
            row += 1


def _add_xlsx_section_sheet(
    wb, block: _RenderBlock, header_fmt, cell_fmt, note_fmt, used: set,
):
    """One worksheet per H2 section. Tables become native cells; prose becomes notes."""
    sheet_name = _xlsx_safe_sheet_name(block.title or "Section", used)
    ws = wb.add_worksheet(sheet_name)
    ws.set_column(0, 0, 28)

    tables, prose_lines = _extract_tables_and_prose(block.body_lines)

    row = 0
    for tbl in tables:
        # Determine column count from the widest row
        cols = max(len(r) for r in tbl)
        # Set reasonable widths beyond column A
        for c in range(1, cols):
            ws.set_column(c, c, 18)

        # Row 0 of the table is the header; subsequent rows are data
        header_row = tbl[0]
        for c, cell in enumerate(header_row):
            ws.write_string(row, c, _strip_inline_md(cell.strip()), header_fmt)
        row += 1
        for data_row in tbl[1:]:
            for c in range(cols):
                cell = data_row[c].strip() if c < len(data_row) else ""
                cell = _strip_inline_md(cell)
                # Try to coerce to number for nice numeric formatting
                num = _try_parse_number(cell)
                if num is not None:
                    ws.write_number(row, c, num, cell_fmt)
                else:
                    ws.write_string(row, c, cell, cell_fmt)
            row += 1
        row += 1  # blank row between tables

    # Prose lines as italic notes in column A
    if prose_lines:
        if tables:
            ws.write_string(row, 0, "Notes", header_fmt)
            row += 1
        for line in prose_lines:
            stripped = line.strip()
            if not stripped:
                continue
            ws.write_string(row, 0, _strip_inline_md(stripped), note_fmt)
            row += 1


def _extract_tables_and_prose(body_lines: list) -> tuple:
    """Walk body lines; return (list_of_tables, list_of_prose_lines).

    Each table is list[list[str]] with the alignment row dropped. Prose lines
    are anything outside tables — including bullets, paragraphs, fenced code
    body, and headings.
    """
    tables: list[list[list[str]]] = []
    current_table: list[list[str]] = []
    prose: list[str] = []
    in_code = False

    def flush_table():
        if not current_table:
            return
        rows = [
            r for r in current_table
            if not all(re.match(r"^[-:\s]+$", c.strip()) for c in r)
        ]
        if rows:
            tables.append(rows)
        current_table.clear()

    for line in body_lines:
        if line.startswith("```"):
            flush_table()
            in_code = not in_code
            continue
        if in_code:
            prose.append(line)
            continue
        if "|" in line and line.strip().startswith("|"):
            cells = [c for c in line.strip().strip("|").split("|")]
            current_table.append(cells)
            continue
        if current_table:
            flush_table()
        prose.append(line)

    flush_table()
    return (tables, prose)


def _try_parse_number(text: str):
    """Attempt to coerce text to a number for native Excel numeric cells.

    Strips $, commas, %; treats trailing % as multiplied by 0.01. Returns None
    on any failure. Conservative — anything that doesn't parse cleanly stays a
    string. Currency or unit context is preserved as text.
    """
    if not text:
        return None
    t = text.strip()
    if not t:
        return None
    # Reject obvious non-numbers fast
    if any(c.isalpha() for c in t):
        return None
    is_pct = t.endswith("%")
    if is_pct:
        t = t[:-1].strip()
    # Strip currency and thousands separators
    t = t.replace("$", "").replace(",", "").replace(" ", "")
    try:
        n = float(t)
    except ValueError:
        return None
    if is_pct:
        n = n / 100.0
    return n


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


# ---------- CLI ----------

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 4:
        print("Usage: python output_runtime.py <off-spec-path> <corpus-instance-path> <output-dir>")
        sys.exit(1)
    result = o_render(sys.argv[1], sys.argv[2], sys.argv[3])
    if result.success:
        print(f"OK: rendered to {result.artifact_path}")
        print(f"Sections read: {result.sections_read}")
    else:
        print(f"FAIL: {result.error}")
        sys.exit(2)
