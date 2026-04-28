"""
format_convert.py — Document Format Conversion Utility (Phase 10, Step 0)

Converts various document formats to clean markdown while preserving structure.
Standalone utility — usable independently for any format conversion need.

Supported formats:
  - PDF (.pdf) — text extraction with heading detection
  - Word (.docx) — text, headings, tables, lists
  - PowerPoint (.pptx) — slide text and speaker notes
  - HTML (.html, .htm) — strip markup, preserve structure
  - RTF (.rtf) — strip formatting
  - Plain text (.txt) — passthrough with minimal cleanup
  - Markdown (.md) — passthrough

Usage:
    from orchestrator.tools.format_convert import convert_to_markdown
    markdown = convert_to_markdown("/path/to/document.pdf")
"""

from __future__ import annotations

import os
import re
from pathlib import Path


def convert_to_markdown(file_path: str) -> str:
    """
    Convert a document to clean markdown.

    Args:
        file_path: Absolute path to the source document.

    Returns:
        Clean markdown string with heading structure preserved.

    Raises:
        ValueError: If the format is not supported.
        FileNotFoundError: If the file doesn't exist.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()

    converters = {
        ".pdf": _convert_pdf,
        ".docx": _convert_docx,
        ".pptx": _convert_pptx,
        ".html": _convert_html,
        ".htm": _convert_html,
        ".rtf": _convert_rtf,
        ".txt": _convert_text,
        ".md": _convert_markdown,
    }

    converter = converters.get(ext)
    if converter is None:
        raise ValueError(f"Unsupported format: {ext}")

    raw = converter(file_path)
    return _cleanup(raw)


def detect_format(file_path: str) -> str:
    """Detect the document format from extension."""
    ext = Path(file_path).suffix.lower()
    format_map = {
        ".pdf": "pdf", ".docx": "docx", ".pptx": "pptx",
        ".html": "html", ".htm": "html", ".rtf": "rtf",
        ".txt": "text", ".md": "markdown",
    }
    return format_map.get(ext, "unknown")


# ---------------------------------------------------------------------------
# Format-specific converters
# ---------------------------------------------------------------------------

def _convert_pdf(file_path: str) -> str:
    """Convert PDF to markdown using pdfplumber."""
    import pdfplumber

    parts = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                parts.append(text)

            # Extract tables
            tables = page.extract_tables()
            for table in tables:
                if table:
                    parts.append(_table_to_markdown(table))

    return "\n\n".join(parts)


def _convert_docx(file_path: str) -> str:
    """Convert Word document to markdown."""
    from docx import Document

    doc = Document(file_path)
    parts = []

    for para in doc.paragraphs:
        style = para.style.name.lower() if para.style else ""
        text = para.text.strip()
        if not text:
            parts.append("")
            continue

        # Map heading styles
        if "heading 1" in style:
            parts.append(f"# {text}")
        elif "heading 2" in style:
            parts.append(f"## {text}")
        elif "heading 3" in style:
            parts.append(f"### {text}")
        elif "heading 4" in style:
            parts.append(f"#### {text}")
        elif "list" in style:
            parts.append(f"- {text}")
        else:
            parts.append(text)

    # Extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            parts.append(_table_to_markdown(rows))

    return "\n\n".join(parts)


def _convert_pptx(file_path: str) -> str:
    """Convert PowerPoint to markdown — slide text and speaker notes."""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)

        if slide_text:
            # First text element as slide heading
            parts.append(f"## Slide {i}: {slide_text[0]}")
            for text in slide_text[1:]:
                parts.append(f"- {text}")

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"\n*Speaker notes:* {notes}")

        parts.append("")  # Blank line between slides

    return "\n\n".join(parts)


def _convert_html(file_path: str) -> str:
    """Convert HTML to markdown using markdownify."""
    from markdownify import markdownify
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()

    # Clean up with BeautifulSoup first
    soup = BeautifulSoup(html, "html.parser")

    # Remove script and style tags
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    # Convert to markdown
    md = markdownify(str(soup), heading_style="ATX", strip=["img"])
    return md


def _convert_rtf(file_path: str) -> str:
    """Convert RTF to plain text."""
    from striprtf.striprtf import rtf_to_text

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        rtf = f.read()

    return rtf_to_text(rtf)


def _convert_text(file_path: str) -> str:
    """Read plain text with minimal cleanup."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _convert_markdown(file_path: str) -> str:
    """Passthrough for markdown files."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _table_to_markdown(rows: list[list]) -> str:
    """Convert a list of rows to a markdown table."""
    if not rows or not rows[0]:
        return ""

    # Header row
    header = rows[0]
    lines = ["| " + " | ".join(str(c) for c in header) + " |"]
    lines.append("|" + "|".join("---" for _ in header) + "|")

    # Data rows
    for row in rows[1:]:
        # Pad or truncate to match header width
        padded = list(row) + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(str(c) for c in padded[:len(header)]) + " |")

    return "\n".join(lines)


def _cleanup(text: str) -> str:
    """Clean up converted text: normalize whitespace, fix common artifacts."""
    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # Collapse 3+ consecutive blank lines to 2
    text = re.sub(r'\n{3,}', '\n\n', text)

    # Remove trailing whitespace from lines
    lines = [line.rstrip() for line in text.split("\n")]
    text = "\n".join(lines)

    # Strip leading/trailing whitespace
    text = text.strip()

    return text


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: format_convert.py <file_path>")
        print("Converts a document to clean markdown and prints to stdout.")
        print(f"Supported formats: {', '.join(['.pdf', '.docx', '.pptx', '.html', '.rtf', '.txt', '.md'])}")
        sys.exit(1)

    file_path = sys.argv[1]
    try:
        result = convert_to_markdown(file_path)
        print(result)
    except (ValueError, FileNotFoundError) as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
